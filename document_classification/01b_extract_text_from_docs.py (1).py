# Databricks notebook source
# MAGIC %pip install pypdf
# MAGIC %pip install pymupdf
# MAGIC %pip install extract-msg
# MAGIC %pip install azure-ai-formrecognizer openai
# MAGIC %pip install presidio-analyzer presidio-anonymizer spacy langdetect azure-ai-formrecognizer pillow
# MAGIC %pip install azure-search-documents
# MAGIC %pip install --upgrade langchain langchain-community
# MAGIC %pip install -U openai
# MAGIC %pip install -U "azure-search-documents"
# MAGIC %pip install mlflow>=3.0 --upgrade
# MAGIC %pip install python-docx
# MAGIC %pip install openpyxl
# MAGIC %pip install xlrd
# MAGIC %pip install python-pptx

# COMMAND ----------

# MAGIC %pip install --pre -U "azure-search-documents>=11.6.0b1"
# MAGIC dbutils.library.restartPython()
# MAGIC

# COMMAND ----------

import sys
sys.path.append("/Workspace/Users/aswin@eyaswin.onmicrosoft.com")

from utils_common import post_json
from map_partitions_api import map_partitions_api
import openai
print("openai version "+ openai.__version__)

import azure.search.documents
print(azure.search.documents.__version__)   # should be 11.7.x or newer



# COMMAND ----------

from pyspark.sql import SparkSession
from pyspark.ml.feature import HashingTF, IDF, Tokenizer, StringIndexer, IndexToString
from pyspark.ml.classification import LogisticRegression
from pyspark.ml import Pipeline
from pyspark.sql import SparkSession
from pyspark.sql.functions import concat_ws, col
from pyspark.sql.functions import regexp_extract

from pyspark.sql import functions as F, types as T
import io, json, csv, email
from pypdf import PdfReader
import extract_msg   # pip install extract-msg
import base64, json, os

# Initialize SparkSession for Databricks
# In Databricks, the SparkSession is typically pre-configured,
# so you don't need to explicitly set Hadoop Azure dependencies.
spark = (
    SparkSession.builder
    .appName("Document-Classification")
    .getOrCreate()
)

configs = {
  "fs.azure.account.auth.type": "OAuth",
  "fs.azure.account.oauth.provider.type": "org.apache.hadoop.fs.azurebfs.oauth2.ClientCredsTokenProvider",
  "fs.azure.account.oauth2.client.id": "640059f9-cebf-47e5-a3f1-7bbe06e7e4a3",
  "fs.azure.account.oauth2.client.secret": "",
  "fs.azure.account.oauth2.client.endpoint": "https://login.microsoftonline.com/115ee48c-5146-4054-9f33-83e2bfe089fd/oauth2/token"
}

spark.conf.set("fs.azure.account.auth.type.docclassifierstoragecct.dfs.core.windows.net", "OAuth")
for key, val in configs.items():
    spark.conf.set(f"{key}.docclassifierstoragecct.dfs.core.windows.net", val)

#configuring spark parallelism
#**Parallelism**: Instead of single-node Python loops for clarity. For more scale, convert the work to Spark parallelism: read the list of file paths to a DataFrame and `mapPartitions` / `foreachPartition` or use UDFs to parallelize extraction and calls.

spark.conf.set("spark.sql.shuffle.partitions", "400")
spark.conf.set("spark.databricks.io.cache.enabled", "true")

#todo
#Large files / memory**: For very large PDFs or huge images you may want streaming approaches. The above copies files locally before processing; adapt if you have memory concerns.
#todo
#**Idempotency**: The code stores SHA256 for each saved image. You can use that to skip reprocessing duplicates.

#account = dbutils.secrets.get("kv-scope","storage_account_name")
account="docclassifierstoragecct"
container_raw   = "raw"
container_stage = "stage"
container_redacted = "redacted"
container_processed="docclassifiercontainer"
abfss = lambda container, path="": f"abfss://{container}@{account}.dfs.core.windows.net/{path}"


import base64, json, os

DOCUMENTINTEL_ENDPOINT = "https://documentsclassifier.cognitiveservices.azure.com/" #dbutils.secrets.get("kv-scope","di_endpoint")  # e.g., https://<res>.cognitiveservices.azure.com
DOCUMENTINTEL_KEY      = ""# dbutils.secrets.get("kv-scope","di_key")
DOCUMENTINTEL_MODEL_ID = "prebuilt-document"  # or your custom model

# --- Azure Configurations ---
AZURE_OPENAI_EMBEDDING_ENDPOINT = "https://tpapp.openai.azure.com/openai/deployments/text-embedding-3-large/embeddings?api-version=2024-02-01"
AZURE_OPENAI_API_KEY = ""
AZURE_SEARCH_ENDPOINT = "https://docsclassifieraisearch.search.windows.net"
AZURE_SEARCH_KEY = ""
AZURE_SEARCH_INDEX = "mydocs-knowledgeharvester-index"

AZURE_OPENAI_API_KEY = ""
AZURE_OPENAI_ENDPOINT = "https://YOUR-AZURE-OPENAI-ENDPOINT.openai.azure.com/"
AZURE_OPENAI_DEPLOYMENT = "gpt-4o-mini"
AZURE_OPENAI_API_VERSION="2024-12-01-preview"


# system folders
PATH_DOCS_IN   = abfss(container_raw,   "incoming/docs/")     # .pdf/.docx/.xlsx will be stored here
PATH_IMAGES_OUT= abfss(container_stage, "images/")            # extracted images go here
PATH_JSON_OUT  = abfss(container_stage, "extracted_json/")    # DI/Vision JSON here
PATH_CLEANED   = abfss(container_redacted, "cleaned/")           # PII-clean text
PATH_EMBED     = abfss(container_stage, "embeddings/")        # parquet with embeddings
PATH_TEXT_OUT = abfss(container_stage, "text/")


# COMMAND ----------

# CSV, eml, msg spark dataframe
from pyspark.sql.functions import monotonically_increasing_id, col
from pyspark.sql import functions as F, types as T
import io, json, csv, email
import extract_msg   # pip install extract-msg
import uuid

PATH_TEXT_OUT = abfss(container_stage, "text/")

# ===========================================================
# 1) Extractors for unsupported formats
# ===========================================================

def extract_text_from_csv(content: bytes) -> str:
    """Return structured JSON string with rows and columns preserved."""
    try:
        with io.StringIO(content.decode(errors="ignore")) as buf:
            reader = csv.DictReader(buf)
            rows = [row for row in reader]
            return json.dumps(rows, ensure_ascii=False)
    except Exception as e:
        return f"[ERROR extracting CSV: {e}]"

def extract_text_from_eml(content: bytes) -> str:
    try:
        msg = email.message_from_bytes(content)
        parts = []
        if msg.is_multipart():
            for part in msg.walk():
                if part.get_content_type() == "text/plain":
                    parts.append(part.get_payload(decode=True).decode(errors="ignore"))
        else:
            parts.append(msg.get_payload(decode=True).decode(errors="ignore"))
        return "\n".join(parts)
    except Exception as e:
        return f"[ERROR extracting EML: {e}]"

def extract_text_from_msg(content: bytes) -> str:
    try:
        with io.BytesIO(content) as buf:
            msg = extract_msg.Message(buf)
            subject = msg.subject or ""
            body = msg.body or ""
            return "\n".join([subject, body])
    except Exception as e:
        return f"[ERROR extracting MSG: {e}]"

# ===========================================================
# 2) Dispatcher
# ===========================================================

def extract_text(row):
    #ensures to access Spark Row fields correctly using attribute access
    path = row.path.lower()
    content = row.content

    if path.endswith(".csv"):
        return extract_text_from_csv(content), "csv"
    if path.endswith(".eml"):
        return extract_text_from_eml(content), "eml"
    if path.endswith(".msg"):
        return extract_text_from_msg(content), "msg"
    if path.endswith((".pdf", ".docx", ".xlsx", ".pptx")):
        return "[UNSUPPORTED FILE TYPE FOR TASK 1B]", "unsupported"
    return "[UNSUPPORTED FILE TYPE FOR TASK 1B]", "unsupported"

# ===========================================================
# 3) Spark Job — iterate over files and extract text
# ===========================================================

files = (spark.read.format("binaryFile")
         .load(PATH_DOCS_IN)   # same input as Task 1
         .select("path","content"))

def map_partition(it):
    for r in it:
        chunk_text, stype = extract_text(r)
        #chunk_text, stype = extract_text(bytes(r["content"]))
        yield {"path": r.path, "chunk_text": chunk_text, "source_type": stype}
#defining the schema for the dataframe
schema = T.StructType([
    T.StructField("path", T.StringType()),
    T.StructField("chunk_text", T.StringType()),
    T.StructField("source_type", T.StringType())
])
print("Number of Partitions ",files.rdd.getNumPartitions())
rdd = files.rdd.mapPartitions(map_partition)
# create dataframe with the schema
# The schema is used to ensure that the DataFrame(JSON) columns have the correct names and types.
#df_text = spark.createDataFrame(rdd, schema)
df_text = spark.createDataFrame(files.rdd.mapPartitions(map_partition), schema)

# repartition DataFrame before writing (e.g., 200 partitions)
# if partition is too small right of task scheduling overhead increases meaning spark has to launch 200 tasks for very tiny chunks of data which will slow down the job. and if its too large then out-of-memory(OOM) may occur. Typiacally between 100 and 200 partitions is a good starting point.
# example 100GB data then 100GB/200=0.5GB per partition which is 512 MB risk of OOM.
#df_text = df_text.repartition(16)
# ===========================================================
# 4) Write results to ADLS
# ===========================================================
#This ensures the folder_name column is available for partitioning when writing the output and writes the output to the folder by the file name.
df_text = (
    df_text
    .withColumn("folder_name", F.regexp_extract("path", r"/([^/]+)\.[^/.]+$", 1))
    .withColumn("unique_id", F.lit(uuid.uuid4().hex))
)

df_text.persist()
print(f"Total rows extracted: {df_text.count()}")

def write_row(row):
    folder = row["folder_name"]
    file_id = row["unique_id"]
    out_path = f"{PATH_TEXT_OUT}/{folder}/{file_id}.parquet"
    single_df = spark.createDataFrame([row.asDict()])
    single_df.write.mode("overwrite").option("compression", "snappy").parquet(out_path)

# Assign random partition IDs
df_text = df_text.withColumn(
    "rand_part", (monotonically_increasing_id() % 1000).cast("int")  # 1000 = desired parallelism
)
output_base = PATH_TEXT_OUT

# Repartition for high parallelism and write with small files
def write_partitioned(batch_df, batch_id):
    # Write each micro-batch (forEachBatch ensures driver-side write planning)
    (batch_df
     .repartition("folder_name", "rand_part")      # create multiple writers per folder
     .write
     .mode("append")                               # append to allow multiple small files per folder
     .option("compression", "snappy")
     .option("maxRecordsPerFile", 1)               # one record per file → ≤250 KB
     .format("parquet")
     .save(output_base)
    )
write_partitioned(df_text, 0)
print("Write complete — each record written as a small Parquet file per folder.")

# Collect the paths to process on the driver
paths_to_move = df_text.select("path").distinct().collect()

# This order guarantees that only files which have been processed and written are moved to the processed container.
for row in paths_to_move:
    src_path = row["path"]
    # Convert abfss:// to dbfs:/mnt/ if you have mounted, or use abfss:// directly if not
    # Here, we replace the container and folder in the path string
    dst_path = src_path.replace(
        f"/{container_raw}/incoming/docs/",
        f"/{container_processed}/processed/docs/"
    )
    dbutils.fs.mv(src_path, dst_path)

# COMMAND ----------

#pdf extraction
from pyspark.sql import functions as F, types as T
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
import io
import nest_asyncio
import itertools
import fitz  # PyMuPDF
import os
import logging

from langchain.text_splitter import RecursiveCharacterTextSplitter

nest_asyncio.apply()

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("pdf_extraction")

files = (
    spark.read.format("binaryFile")
    .option("recursiveFileLookup", "true")
    .load(PATH_DOCS_IN)
    .filter(F.lower(F.col("path")).endswith(".pdf"))
    .select("path", "content")
)

schema = T.StructType([
    T.StructField("path", T.StringType()),
    T.StructField("chunk_index", T.IntegerType()),
    T.StructField("chunk_text", T.StringType()),
    T.StructField("page_spans", T.ArrayType(T.IntegerType())),
    T.StructField("source_type", T.StringType()),
    T.StructField("table_json", T.StringType()),
    T.StructField("image_count", T.IntegerType()),
    T.StructField("image_info", T.ArrayType(
        T.StructType([
            T.StructField("page_num", T.IntegerType()),
            T.StructField("img_idx", T.IntegerType()),
            T.StructField("image_bytes", T.BinaryType())
        ])
    ))
])

# Counter for Document Intelligence calls
docintel_call_counter = {}

def analyze_pdf_with_azure(content_bytes, doc_path):
    endpoint = DOCUMENTINTEL_ENDPOINT
    key = DOCUMENTINTEL_KEY
    client = DocumentAnalysisClient(endpoint=endpoint, credential=AzureKeyCredential(key))
    poller = client.begin_analyze_document(DOCUMENTINTEL_MODEL_ID, document=io.BytesIO(content_bytes))
    # Log the call
    docintel_call_counter[doc_path] = docintel_call_counter.get(doc_path, 0) + 1
    logger.info(f"Document Intelligence called for: {doc_path} (Total calls: {docintel_call_counter[doc_path]})")
    return poller.result()
# Parses tables from the Azure result.
def extract_tables_from_azure(result):
    tables = []
    for table in getattr(result, "tables", []):
        table_data = {
            "row_count": table.row_count,
            "column_count": table.column_count,
            "cells": [
                {
                    "row_index": cell.row_index,
                    "column_index": cell.column_index,
                    "content": cell.content
                }
                for cell in table.cells
            ]
        }
        tables.append(table_data)
    return tables
# Use PyMuPDF (fitz) to extract text and images from each page of the PDF.
def extract_text_and_images_with_fitz(content_bytes, doc_path):
    doc = fitz.open(stream=content_bytes, filetype="pdf")
    text_chunks = []
    image_count = 0
    image_info = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        text_chunks.append(page.get_text())
        images = page.get_images(full=True)
        image_count += len(images)
        for img_idx, img in enumerate(images):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_info.append({
                "page_num": page_num + 1,
                "img_idx": img_idx + 1,
                "image_bytes": image_bytes
            })
    return text_chunks, image_count, image_info

# Chunking function using LangChain's RecursiveCharacterTextSplitter
def chunk_text_with_overlap(text, chunk_size=1000, chunk_overlap=150):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ".", " ", ""]
    )
    return splitter.split_text(text)

# Combines the extracted tables, text and images yielding a record for each chunk with text, table data, and image info.
def extract_chunks_from_pdf(row):
    try:
        content_bytes = row["content"]
        doc_path = row["path"]
        # Extract structured tables using Azure Document Intelligence
        azure_result = analyze_pdf_with_azure(content_bytes, doc_path)
        tables = extract_tables_from_azure(azure_result)
        table_json = json.dumps(tables, ensure_ascii=False)
        # Extract text and images using PyMuPDF
        text_pages, image_count, image_info = extract_text_and_images_with_fitz(content_bytes, doc_path)
        chunk_idx = 0
        for page_num, page_text in enumerate(text_pages):
            # Split each page's text into overlapping chunks
            page_chunks = chunk_text_with_overlap(page_text)
            for local_idx, chunk in enumerate(page_chunks):
                yield {
                    "path": doc_path,
                    "chunk_index": chunk_idx,
                    "chunk_text": chunk,
                    "page_spans": [page_num + 1],
                    "source_type": "pdf",
                    "table_json": table_json if chunk_idx == 0 else None,
                    "image_count": image_count if chunk_idx == 0 else None,
                    "image_info": image_info if chunk_idx == 0 else None
                }
                chunk_idx += 1
    except Exception as e:
        logger.error(f"[ERROR extracting pdf: {e}] for document {row.get('path', 'unknown')}")
        yield {
            "path": row["path"],
            "chunk_index": -1,
            "chunk_text": f"[ERROR extracting pdf: {e}]",
            "page_spans": [],
            "source_type": "pdf_error",
            "table_json": None,
            "image_count": None,
            "image_info": None
        }

def map_partitions_api_word(rows, func):
    for row in rows:
        if hasattr(row, "asDict"):
            yield from func(row.asDict())
        else:
            yield from func(row)
# using Spark RDD's mapPartitions to process PDFs in parallel so each partition processes a batch of PDFs and extracts their content.
def partition_mapper(rows):
    return map_partitions_api_word(rows, extract_chunks_from_pdf)

rdd = files.rdd.mapPartitions(
    lambda rows: itertools.chain.from_iterable([partition_mapper(rows)])
)
# create a spark dataframe df_chunks from the extracted records using the schema defined earlier 
df_chunks = spark.createDataFrame(rdd, schema)
# repartitions dataframe
#df_chunks = df_chunks.repartition(16)
#  Add a folder_name column for partitioning
df_chunks = df_chunks.withColumn(
    "folder_name",
    F.regexp_extract("path", r"/([^/]+)\.[^/.]+$", 1)
)
# Write the extracted data as Parquet files to ADLS PATH_TEXT_OUT partitioned by folder_name
(df_chunks.write
    .mode("overwrite")
    .partitionBy("folder_name")
    .parquet(PATH_TEXT_OUT))

# Collect image info to driver and save images to ADLS using dbutils
image_rows = (
    df_chunks
    .filter(F.col("image_info").isNotNull())
    .select("path", "image_info")
    .collect()
)

for row in image_rows:
    doc_path = row["path"]
    image_info_list = row["image_info"]
    base_name = os.path.splitext(os.path.basename(doc_path))[0]
    for img in image_info_list:
        page_num = img["page_num"]
        img_idx = img["img_idx"]
        image_bytes = img["image_bytes"]
        folder_path = f"{PATH_IMAGES_OUT}/{base_name}/page_{page_num}"
        image_name = f"image_{img_idx}.png"
        full_path = f"{folder_path}/{image_name}"
        dbutils.fs.mkdirs(folder_path)
        with open("/tmp/tmp_img.png", "wb") as f:
            f.write(image_bytes)
        dbutils.fs.cp(f"file:/tmp/tmp_img.png", full_path, True)

display(df_chunks.limit(10))

# COMMAND ----------

#word extraction
from pyspark.sql import functions as F, types as T
import io
import nest_asyncio
import itertools
import os
import logging
import json
from langchain.text_splitter import RecursiveCharacterTextSplitter
from docx import Document  # python-docx

nest_asyncio.apply()

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("word_extraction")

files = (
    spark.read.format("binaryFile")
    .option("recursiveFileLookup", "true")
    .load(PATH_DOCS_IN)
    .filter(F.lower(F.col("path")).endswith(".docx"))
    .select("path", "content")
)

schema = T.StructType([
    T.StructField("path", T.StringType()),
    T.StructField("chunk_index", T.IntegerType()),
    T.StructField("chunk_text", T.StringType()),
    T.StructField("page_spans", T.ArrayType(T.IntegerType())),
    T.StructField("source_type", T.StringType()),
    T.StructField("table_json", T.StringType()),
    T.StructField("image_count", T.IntegerType()),
    T.StructField("image_info", T.ArrayType(
        T.StructType([
            T.StructField("page_num", T.IntegerType()),
            T.StructField("img_idx", T.IntegerType()),
            T.StructField("image_bytes", T.BinaryType())
        ])
    ))
])

def extract_tables_from_docx(doc):
    tables = []
    for table in doc.tables:
        table_data = {
            "row_count": len(table.rows),
            "column_count": len(table.columns),
            "cells": [
                {
                    "row_index": i,
                    "column_index": j,
                    "content": cell.text
                }
                for i, row in enumerate(table.rows)
                for j, cell in enumerate(row.cells)
            ]
        }
        tables.append(table_data)
    return tables

def extract_text_and_images_from_docx(content_bytes, doc_path):
    text_chunks = []
    image_info = []
    image_count = 0
    tables = []
    try:
        with io.BytesIO(content_bytes) as buf:
            doc = Document(buf)
            # Extract text
            for para in doc.paragraphs:
                text_chunks.append(para.text)
            # Extract tables as text and as structured data
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    table_text.append("\t".join(row_text))
                text_chunks.append("\n".join(table_text))
            tables = extract_tables_from_docx(doc)
            # Extract images
            rels = doc.part.rels
            img_idx = 1
            for rel in rels:
                rel_obj = rels[rel]
                if "image" in rel_obj.target_ref:
                    image_bytes = rel_obj.target_part.blob
                    image_info.append({
                        "page_num": 1,
                        "img_idx": img_idx,
                        "image_bytes": image_bytes
                    })
                    img_idx += 1
            image_count = len(image_info)
    except Exception as e:
        logger.error(f"[ERROR extracting docx: {e}] for document {doc_path}")
    return text_chunks, image_count, image_info, tables

def chunk_text_with_overlap(text, chunk_size=1000, chunk_overlap=150):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ".", " ", ""]
    )
    return splitter.split_text(text)

def extract_chunks_from_docx(row):
    try:
        content_bytes = row["content"]
        doc_path = row["path"]
        text_chunks, image_count, image_info, tables = extract_text_and_images_from_docx(content_bytes, doc_path)
        table_json = json.dumps(tables, ensure_ascii=False)
        full_text = "\n".join(text_chunks)
        page_chunks = chunk_text_with_overlap(full_text)
        for chunk_idx, chunk in enumerate(page_chunks):
            yield {
                "path": doc_path,
                "chunk_index": chunk_idx,
                "chunk_text": chunk,
                "page_spans": [],
                "source_type": "docx",
                "table_json": table_json if chunk_idx == 0 else None,
                "image_count": image_count if chunk_idx == 0 else None,
                "image_info": image_info if chunk_idx == 0 else None
            }
    except Exception as e:
        logger.error(f"[ERROR extracting docx: {e}] for document {row.get('path', 'unknown')}")
        yield {
            "path": row["path"],
            "chunk_index": -1,
            "chunk_text": f"[ERROR extracting docx: {e}]",
            "page_spans": [],
            "source_type": "docx_error",
            "table_json": None,
            "image_count": None,
            "image_info": None
        }

def map_partitions_api_word(rows, func):
    for row in rows:
        if hasattr(row, "asDict"):
            yield from func(row.asDict())
        else:
            yield from func(row)

def partition_mapper(rows):
    return map_partitions_api_word(rows, extract_chunks_from_docx)

rdd = files.rdd.mapPartitions(
    lambda rows: itertools.chain.from_iterable([partition_mapper(rows)])
)
df_chunks = spark.createDataFrame(rdd, schema)

# Add folder_name column before writing
df_chunks = df_chunks.withColumn(
    "folder_name",
    F.regexp_extract("path", r"/([^/]+)\.[^/.]+$", 1)
)

# Dynamically choose format and partition count based on dataset size
record_count = df_chunks.count()
if record_count < 10000:
    num_partitions = 4
elif record_count < 100000:
    num_partitions = 16
else:
    num_partitions = 64

use_delta = False  # Set to False to force Parquet

if use_delta:
    df_chunks = df_chunks.repartition(num_partitions)
    (df_chunks.write
        .mode("overwrite")
        .format("delta")
        .partitionBy("folder_name")
        .save(PATH_TEXT_OUT))
else:
    df_chunks = df_chunks.repartition(num_partitions)
    (df_chunks.write
        .mode("overwrite")
        .partitionBy("folder_name")
        .parquet(PATH_TEXT_OUT))

image_rows = (
    df_chunks
    .filter(F.col("image_info").isNotNull())
    .select("path", "image_info")
    .collect()
)

for row in image_rows:
    doc_path = row["path"]
    image_info_list = row["image_info"]
    base_name = os.path.splitext(os.path.basename(doc_path))[0]
    for img in image_info_list:
        page_num = img["page_num"]
        img_idx = img["img_idx"]
        image_bytes = img["image_bytes"]
        folder_path = f"{PATH_IMAGES_OUT}/{base_name}/page_{page_num}"
        image_name = f"image_{img_idx}.png"
        full_path = f"{folder_path}/{image_name}"
        dbutils.fs.mkdirs(folder_path)
        with open("/tmp/tmp_img.png", "wb") as f:
            f.write(image_bytes)
        dbutils.fs.cp(f"file:/tmp/tmp_img.png", full_path, True)

display(df_chunks.limit(10))

# COMMAND ----------

# xlsx/xls extraction
from pyspark.sql import functions as F, types as T
import io
import nest_asyncio
import itertools
import os
import logging
import json
from langchain.text_splitter import RecursiveCharacterTextSplitter
import openpyxl  # for .xlsx, .xlsm
import xlrd      # for .xls
from PIL import Image
from openpyxl.drawing.image import Image as OpenpyxlImage

nest_asyncio.apply()

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("excel_extraction")

files = (
    spark.read.format("binaryFile")
    .option("recursiveFileLookup", "true")
    .load(PATH_DOCS_IN)
    .filter(
        (F.lower(F.col("path")).endswith(".xlsx")) |
        (F.lower(F.col("path")).endswith(".xls")) |
        (F.lower(F.col("path")).endswith(".xlsm"))
    )
    .select("path", "content")
)

# Change page_spans to ArrayType(StringType()) to avoid Delta merge issues
schema = T.StructType([
    T.StructField("path", T.StringType()),
    T.StructField("chunk_index", T.IntegerType()),
    T.StructField("chunk_text", T.StringType()),
    T.StructField("page_spans", T.ArrayType(T.StringType())),
    T.StructField("source_type", T.StringType()),
    T.StructField("table_json", T.StringType()),
    T.StructField("image_count", T.IntegerType()),
    T.StructField("image_info", T.ArrayType(
        T.StructType([
            T.StructField("sheet_name", T.StringType()),
            T.StructField("img_idx", T.IntegerType()),
            T.StructField("image_bytes", T.BinaryType())
        ])
    ))
])

def extract_tables_from_openpyxl(ws):
    rows = list(ws.iter_rows(values_only=True))
    if not rows:
        return []
    table_data = {
        "row_count": len(rows),
        "column_count": len(rows[0]) if rows else 0,
        "cells": [
            {
                "row_index": i,
                "column_index": j,
                "content": str(cell) if cell is not None else ""
            }
            for i, row in enumerate(rows)
            for j, cell in enumerate(row)
        ]
    }
    return [table_data]

def extract_tables_from_xlrd(sheet):
    rows = [sheet.row_values(i) for i in range(sheet.nrows)]
    if not rows:
        return []
    table_data = {
        "row_count": len(rows),
        "column_count": len(rows[0]) if rows else 0,
        "cells": [
            {
                "row_index": i,
                "column_index": j,
                "content": str(cell) if cell is not None else ""
            }
            for i, row in enumerate(rows)
            for j, cell in enumerate(row)
        ]
    }
    return [table_data]

def extract_text_and_images_from_excel(content_bytes, doc_path):
    text_chunks = []
    image_info = []
    image_count = 0
    tables = []
    page_spans = []
    try:
        ext = doc_path.lower().split(".")[-1]
        if ext in ["xlsx", "xlsm"]:
            with io.BytesIO(content_bytes) as buf:
                wb = openpyxl.load_workbook(buf, data_only=True, keep_vba=True)
                img_idx = 1
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    # Extract text
                    rows = list(ws.iter_rows(values_only=True))
                    for row in rows:
                        text_line = "\t".join([str(cell) if cell is not None else "" for cell in row])
                        text_chunks.append(f"[{sheet_name}] {text_line}")
                    # Extract tables
                    tables.extend(extract_tables_from_openpyxl(ws))
                    # Extract images
                    if hasattr(ws, '_images'):
                        for img in ws._images:
                            if isinstance(img, OpenpyxlImage):
                                img_bytes = io.BytesIO()
                                img.image.save(img_bytes, format="PNG")
                                image_info.append({
                                    "sheet_name": sheet_name,
                                    "img_idx": img_idx,
                                    "image_bytes": img_bytes.getvalue()
                                })
                                img_idx += 1
                    page_spans.append(str(sheet_name))
                image_count = len(image_info)
        elif ext == "xls":
            with io.BytesIO(content_bytes) as buf:
                wb = xlrd.open_workbook(file_contents=buf.read())
                for sheet in wb.sheets():
                    # Extract text
                    for row_idx in range(sheet.nrows):
                        row = sheet.row_values(row_idx)
                        text_line = "\t".join([str(cell) if cell is not None else "" for cell in row])
                        text_chunks.append(f"[{sheet.name}] {text_line}")
                    # Extract tables
                    tables.extend(extract_tables_from_xlrd(sheet))
                    page_spans.append(str(sheet.name))
            # xlrd does not support images in .xls
    except Exception as e:
        logger.error(f"[ERROR extracting excel: {e}] for document {doc_path}")
    return text_chunks, image_count, image_info, tables, page_spans

def chunk_text_with_overlap(text, chunk_size=1000, chunk_overlap=150):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ".", " ", ""]
    )
    return splitter.split_text(text)

def extract_chunks_from_excel(row):
    try:
        content_bytes = row["content"]
        doc_path = row["path"]
        text_chunks, image_count, image_info, tables, page_spans = extract_text_and_images_from_excel(content_bytes, doc_path)
        table_json = json.dumps(tables, ensure_ascii=False)
        full_text = "\n".join(text_chunks)
        page_spans_str = [str(s) for s in page_spans]  # Ensure all elements are string
        page_chunks = chunk_text_with_overlap(full_text)
        for chunk_idx, chunk in enumerate(page_chunks):
            yield {
                "path": doc_path,
                "chunk_index": chunk_idx,
                "chunk_text": chunk,
                "page_spans": page_spans_str,
                "source_type": "excel",
                "table_json": table_json if chunk_idx == 0 else None,
                "image_count": image_count if chunk_idx == 0 else None,
                "image_info": image_info if chunk_idx == 0 else None
            }
    except Exception as e:
        logger.error(f"[ERROR extracting excel: {e}] for document {row.get('path', 'unknown')}")
        yield {
            "path": row["path"],
            "chunk_index": -1,
            "chunk_text": f"[ERROR extracting excel: {e}]",
            "page_spans": [""],  # Ensure this is always a list of strings
            "source_type": "excel_error",
            "table_json": None,
            "image_count": None,
            "image_info": None
        }

def map_partitions_api_excel(rows, func):
    for row in rows:
        if hasattr(row, "asDict"):
            yield from func(row.asDict())
        else:
            yield from func(row)

def partition_mapper(rows):
    return map_partitions_api_excel(rows, extract_chunks_from_excel)

rdd = files.rdd.mapPartitions(
    lambda rows: itertools.chain.from_iterable([partition_mapper(rows)])
)
df_chunks = spark.createDataFrame(rdd, schema)
record_count = df_chunks.count()
if record_count < 10000:
    num_partitions = 4
elif record_count < 100000:
    num_partitions = 16
else:
    num_partitions = 64

df_chunks = df_chunks.withColumn(
    "folder_name",
    F.regexp_extract("path", r"/([^/]+)\.[^/.]+$", 1)
)
df_chunks = df_chunks.repartition(num_partitions)
# Always write as Parquet and do not use Delta
(df_chunks.write
    .mode("overwrite")
    .partitionBy("folder_name")
    .format("parquet")
    .save(PATH_TEXT_OUT)
)
image_rows = (
    df_chunks
    .filter(F.col("image_info").isNotNull())
    .select("path", "image_info")
    .collect()
)

for row in image_rows:
    doc_path = row["path"]
    image_info_list = row["image_info"]
    base_name = os.path.splitext(os.path.basename(doc_path))[0]
    for img in image_info_list:
        sheet_name = img["sheet_name"]
        img_idx = img["img_idx"]
        image_bytes = img["image_bytes"]
        folder_path = f"{PATH_IMAGES_OUT}/{base_name}/{sheet_name}"
        image_name = f"image_{img_idx}.png"
        full_path = f"{folder_path}/{image_name}"
        dbutils.fs.mkdirs(folder_path)
        with open("/tmp/tmp_img.png", "wb") as f:
            f.write(image_bytes)
        dbutils.fs.cp(f"file:/tmp/tmp_img.png", full_path, True)

display(df_chunks.limit(10))

# COMMAND ----------

# pptx/ppt extraction
from pyspark.sql import functions as F, types as T
import io
import nest_asyncio
import itertools
import os
import logging
import json
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pptx import Presentation  # python-pptx
from PIL import Image

nest_asyncio.apply()

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("ppt_extraction")

files = (
    spark.read.format("binaryFile")
    .option("recursiveFileLookup", "true")
    .load(PATH_DOCS_IN)
    .filter(
        (F.lower(F.col("path")).endswith(".pptx")) |
        (F.lower(F.col("path")).endswith(".ppt"))
    )
    .select("path", "content")
)

schema = T.StructType([
    T.StructField("path", T.StringType()),
    T.StructField("chunk_index", T.IntegerType()),
    T.StructField("chunk_text", T.StringType()),
    T.StructField("page_spans", T.ArrayType(T.IntegerType())),
    T.StructField("source_type", T.StringType()),
    T.StructField("table_json", T.StringType()),
    T.StructField("image_count", T.IntegerType()),
    T.StructField("image_info", T.ArrayType(
        T.StructType([
            T.StructField("slide_num", T.IntegerType()),
            T.StructField("img_idx", T.IntegerType()),
            T.StructField("image_bytes", T.BinaryType())
        ])
    ))
])

def extract_text_from_shape(shape):
    """Extract all text from a shape, including text frames and tables."""
    text = ""
    if hasattr(shape, "text") and shape.has_text_frame:
        text = shape.text
    elif shape.shape_type == 19:  # TABLE
        table = shape.table
        rows = []
        for row in table.rows:
            cells = [cell.text for cell in row.cells]
            rows.append("\t".join(cells))
        text = "\n".join(rows)
    return text

def extract_text_and_images_from_ppt(content_bytes, doc_path):
    text_chunks = []
    image_info = []
    image_count = 0
    tables = []
    page_spans = []
    try:
        ext = doc_path.lower().split(".")[-1]
        if ext == "pptx":
            with io.BytesIO(content_bytes) as buf:
                prs = Presentation(buf)
                img_idx = 1
                for slide_num, slide in enumerate(prs.slides, start=1):
                    slide_texts = []
                    for shape in slide.shapes:
                        # Extract text from text frames and tables
                        extracted_text = extract_text_from_shape(shape)
                        if extracted_text:
                            slide_texts.append(extracted_text)
                        # Extract images
                        if shape.shape_type == 13:  # PICTURE
                            image = shape.image
                            image_bytes = image.blob
                            image_info.append({
                                "slide_num": slide_num,
                                "img_idx": img_idx,
                                "image_bytes": image_bytes
                            })
                            img_idx += 1
                    slide_text = f"[Slide {slide_num}] " + "\n".join(slide_texts) if slide_texts else f"[Slide {slide_num}]"
                    text_chunks.append(slide_text)
                    page_spans.append(slide_num)
                image_count = len(image_info)
        elif ext == "ppt":
            text_chunks.append("[UNSUPPORTED: .ppt legacy format. Only .pptx fully supported.]")
    except Exception as e:
        logger.error(f"[ERROR extracting ppt: {e}] for document {doc_path}")
    return text_chunks, image_count, image_info, tables, page_spans

def chunk_text_with_overlap(text, chunk_size=1000, chunk_overlap=150):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ".", " ", ""]
    )
    return splitter.split_text(text)

def extract_chunks_from_ppt(row):
    try:
        content_bytes = row["content"]
        doc_path = row["path"]
        text_chunks, image_count, image_info, tables, page_spans = extract_text_and_images_from_ppt(content_bytes, doc_path)
        table_json = json.dumps(tables, ensure_ascii=False)
        full_text = "\n".join(text_chunks)
        page_chunks = chunk_text_with_overlap(full_text)
        for chunk_idx, chunk in enumerate(page_chunks):
            yield {
                "path": doc_path,
                "chunk_index": chunk_idx,
                "chunk_text": chunk,
                "page_spans": page_spans,
                "source_type": "pptx" if doc_path.lower().endswith(".pptx") else "ppt",
                "table_json": table_json if chunk_idx == 0 else None,
                "image_count": image_count if chunk_idx == 0 else None,
                "image_info": image_info if chunk_idx == 0 else None
            }
    except Exception as e:
        logger.error(f"[ERROR extracting ppt: {e}] for document {row.get('path', 'unknown')}")
        yield {
            "path": row["path"],
            "chunk_index": -1,
            "chunk_text": f"[ERROR extracting ppt: {e}]",
            "page_spans": [],
            "source_type": "ppt_error",
            "table_json": None,
            "image_count": None,
            "image_info": None
        }

def map_partitions_api_ppt(rows, func):
    for row in rows:
        if hasattr(row, "asDict"):
            yield from func(row.asDict())
        else:
            yield from func(row)

def partition_mapper(rows):
    return map_partitions_api_ppt(rows, extract_chunks_from_ppt)

rdd = files.rdd.mapPartitions(
    lambda rows: itertools.chain.from_iterable([partition_mapper(rows)])
)
df_chunks = spark.createDataFrame(rdd, schema)
#df_chunks = df_chunks.repartition(16)
df_chunks = df_chunks.withColumn(
    "folder_name",
    F.regexp_extract("path", r"/([^/]+)\.[^/.]+$", 1)
)
(df_chunks.write
    .mode("overwrite")
    .partitionBy("folder_name")
    .parquet(PATH_TEXT_OUT))

image_rows = (
    df_chunks
    .filter(F.col("image_info").isNotNull())
    .select("path", "image_info")
    .collect()
)

for row in image_rows:
    doc_path = row["path"]
    image_info_list = row["image_info"]
    base_name = os.path.splitext(os.path.basename(doc_path))[0]
    for img in image_info_list:
        slide_num = img["slide_num"]
        img_idx = img["img_idx"]
        image_bytes = img["image_bytes"]
        folder_path = f"/dbfs{PATH_IMAGES_OUT}/{base_name}/slide_{slide_num}"
        image_name = f"image_{img_idx}.png"
        os.makedirs(folder_path, exist_ok=True)
        full_path = f"{folder_path}/{image_name}"
        with open(full_path, "wb") as f:
            f.write(image_bytes)

display(df_chunks.limit(10))

# COMMAND ----------

# msg file extraction
# ===========================================================
# MSG FILE EXTRACTION  (Enhanced: Per-Email Hashing)
# ===========================================================
from pyspark.sql import functions as F, types as T
import io
import nest_asyncio
import itertools
import os
import logging
import json
import hashlib
from langchain.text_splitter import RecursiveCharacterTextSplitter
import extract_msg  # for .msg files

nest_asyncio.apply()

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("msg_extraction")

files = (
    spark.read.format("binaryFile")
    .option("recursiveFileLookup", "true")
    .load(PATH_DOCS_IN)
    .filter(F.lower(F.col("path")).endswith(".msg"))
    .select("path", "content")
)

schema = T.StructType([
    T.StructField("path", T.StringType()),
    T.StructField("chunk_index", T.IntegerType()),
    T.StructField("chunk_text", T.StringType()),
    T.StructField("page_spans", T.ArrayType(T.StringType())),
    T.StructField("source_type", T.StringType()),
    T.StructField("table_json", T.StringType()),
    T.StructField("image_count", T.IntegerType()),
    T.StructField("image_info", T.ArrayType(
        T.StructType([
            T.StructField("attachment_name", T.StringType()),
            T.StructField("img_idx", T.IntegerType()),
            T.StructField("image_bytes", T.BinaryType())
        ])
    )),
    #new metadata fields per updated logic on 10/13/2025
    T.StructField("email_chain_json", T.StringType()),   # structured chain for auditing
    T.StructField("email_hashes", T.ArrayType(T.StringType()))  # hashes for deduplication
])

# ===========================================================
# Utility helpers
# ===========================================================

def compute_hash(content: str) -> str:
    return hashlib.sha256(content.encode("utf-8", errors="ignore")).hexdigest()

def extract_individual_emails(content_bytes: bytes) -> list[dict]:
    """Extract all individual emails (and sub-messages) from a .msg chain."""
    try:
        msg = extract_msg.Message(io.BytesIO(content_bytes))
        chain = []
        messages = [msg] + list(msg.walk()) if hasattr(msg, "walk") else [msg]
        for m in messages:
            subject = getattr(m, "subject", "") or ""
            sender = getattr(m, "sender", "") or ""
            to = getattr(m, "to", "") or ""
            date = str(getattr(m, "date", "")) or ""
            body = getattr(m, "body", "") or ""
            email_hash = compute_hash(body)
            chain.append({
                "subject": subject.strip(),
                "sender": sender.strip(),
                "to": to.strip(),
                "date": date.strip(),
                "body": body.strip(),
                "email_hash": email_hash
            })
        return chain
    except Exception as e:
        logger.error(f"[ERROR extracting individual emails: {e}]")
        return []

def extract_text_and_images_from_msg(content_bytes, doc_path):
    text_chunks, image_info, tables, page_spans = [], [], [], []
    image_count = 0
    try:
        emails = extract_individual_emails(content_bytes)

        # Flatten all bodies for chunking
        full_body_texts = [f"Email {i+1}: {e['body']}" for i, e in enumerate(emails)]
        all_text = "\n\n".join(full_body_texts)

        # optional summary header
        text_chunks.append(f"Extracted {len(emails)} emails from chain at {doc_path}")
        text_chunks.append(all_text)

        # Extract images (attachments from the top-level message)
        msg = extract_msg.Message(io.BytesIO(content_bytes))
        img_idx = 1
        for att in msg.attachments:
            att_name = att.longFilename or att.shortFilename or f"attachment_{img_idx}"
            att_data = att.data
            if att_name.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".bmp")):
                image_info.append({
                    "attachment_name": att_name,
                    "img_idx": img_idx,
                    "image_bytes": att_data
                })
                img_idx += 1
        image_count = len(image_info)
        page_spans.append("msg")

        # New fields for deduplication
        email_hashes = [e["email_hash"] for e in emails]
        email_chain_json = json.dumps(emails, ensure_ascii=False)

        return text_chunks, image_count, image_info, tables, page_spans, email_chain_json, email_hashes

    except Exception as e:
        logger.error(f"[ERROR extracting msg: {e}] for document {doc_path}")
        return [f"[ERROR extracting msg: {e}]"], 0, [], [], [], json.dumps([]), []

def chunk_text_with_overlap(text, chunk_size=1000, chunk_overlap=150):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ".", " ", ""]
    )
    return splitter.split_text(text)

def extract_chunks_from_msg(row):
    try:
        content_bytes = row["content"]
        doc_path = row["path"]

        text_chunks, image_count, image_info, tables, page_spans, email_chain_json, email_hashes = \
            extract_text_and_images_from_msg(content_bytes, doc_path)

        table_json = json.dumps(tables, ensure_ascii=False)
        full_text = "\n".join(text_chunks)
        page_chunks = chunk_text_with_overlap(full_text)

        for chunk_idx, chunk in enumerate(page_chunks):
            yield {
                "path": doc_path,
                "chunk_index": chunk_idx,
                "chunk_text": chunk,
                "page_spans": page_spans,
                "source_type": "msg",
                "table_json": table_json if chunk_idx == 0 else None,
                "image_count": image_count if chunk_idx == 0 else None,
                "image_info": image_info if chunk_idx == 0 else None,
                "email_chain_json": email_chain_json if chunk_idx == 0 else None,
                "email_hashes": email_hashes if chunk_idx == 0 else None
            }
    except Exception as e:
        logger.error(f"[ERROR extracting msg: {e}] for document {row.get('path', 'unknown')}")
        yield {
            "path": row["path"],
            "chunk_index": -1,
            "chunk_text": f"[ERROR extracting msg: {e}]",
            "page_spans": [],
            "source_type": "msg_error",
            "table_json": None,
            "image_count": None,
            "image_info": None,
            "email_chain_json": json.dumps([]),
            "email_hashes": []
        }

def map_partitions_api_msg(rows, func):
    for row in rows:
        if hasattr(row, "asDict"):
            yield from func(row.asDict())
        else:
            yield from func(row)

def partition_mapper(rows):
    return map_partitions_api_msg(rows, extract_chunks_from_msg)

rdd = files.rdd.mapPartitions(
    lambda rows: itertools.chain.from_iterable([partition_mapper(rows)])
)
df_chunks = spark.createDataFrame(rdd, schema)

df_chunks = df_chunks.withColumn(
    "folder_name",
    F.regexp_extract("path", r"/([^/]+)\.[^/.]+$", 1)
)

(df_chunks.write
    .mode("overwrite")
    .partitionBy("folder_name")
    .parquet(PATH_TEXT_OUT))

# ===========================================================
# IMAGE EXTRACTION
# ===========================================================
image_rows = (
    df_chunks
    .filter(F.col("image_info").isNotNull())
    .select("path", "image_info")
    .collect()
)

for row in image_rows:
    doc_path = row["path"]
    image_info_list = row["image_info"]
    base_name = os.path.splitext(os.path.basename(doc_path))[0]
    for img in image_info_list:
        att_name = img["attachment_name"]
        img_idx = img["img_idx"]
        image_bytes = img["image_bytes"]
        folder_path = f"{PATH_IMAGES_OUT}/{base_name}/attachment"
        image_name = f"{att_name or 'image'}_{img_idx}.png"
        full_path = f"{folder_path}/{image_name}"
        dbutils.fs.mkdirs(folder_path)
        with open("/tmp/tmp_img.png", "wb") as f:
            f.write(image_bytes)
        dbutils.fs.cp("file:/tmp/tmp_img.png", full_path, True)

display(df_chunks.limit(10))


# COMMAND ----------

import spacy.cli
spacy.cli.download("xx_ent_wiki_sm")
nlp = spacy.load("xx_ent_wiki_sm")
dbutils.library.restartPython()

# COMMAND ----------

#Presidio detects and redacts PII from each text chunk, and the redacted text is saved in the output DataFrame and written to storage.
import os
import io
import json
import time
from pyspark.sql import functions as F, types as T
from presidio_analyzer import AnalyzerEngine, RecognizerRegistry, PatternRecognizer, Pattern
from presidio_anonymizer import AnonymizerEngine
from langdetect import detect
import spacy
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
from azure.core.exceptions import HttpResponseError
from PIL import Image, ImageDraw, ImageFont


analyzer = AnalyzerEngine(
    nlp_engine=None,
    supported_languages=["en", "nl", "fr", "ko", "hi", "zh", "es", "de","pl", "it"]
)
anonymizer = AnonymizerEngine()

# Add custom EIN recognizer to Presidio
ein_pattern = Pattern("EIN Regex", r"\b\d{2}-\d{7}\b", 0.9)
ein_recognizer = PatternRecognizer(
    supported_entity=["US_EMPLOYER_ID_NUMBER"],
    patterns=[ein_pattern],
    name="ein_recognizer"
)
analyzer.registry.add_recognizer(ein_recognizer)

def detect_language(text):
    try:
        return detect(text)
    except Exception:
        return "en"

detect_language_udf = F.udf(detect_language, T.StringType())

def redact_text(text, lang):
    try:
        results = analyzer.analyze(text=text, language=lang, entities=None)
        if not results:
            return text
        anonymized = anonymizer.anonymize(text=text, analyzer_results=results)
        return anonymized.text
    except Exception:
        return text

redact_text_udf = F.udf(redact_text, T.StringType())

df_text = spark.read.parquet(PATH_TEXT_OUT)
#Spark pipeline to redact PII. The detect_language_udf and redact_text_udf are applied to the DataFrame
df_cleaned = (
    df_text
    .withColumn("language", detect_language_udf("chunk_text"))
    .withColumn("redacted_text", redact_text_udf("chunk_text", "language"))
    .withColumn("folder_name", F.regexp_extract("path", r"/([^/]+)\.[^/.]+$", 1))
)
# Presidio detects and redacts PII from each text chunk, and the redacted text is saved in the output DataFrame and written to storage.
(df_cleaned
    .select("path", "redacted_text", "folder_name")
    .write
    .mode("overwrite")
    .partitionBy("folder_name")
    .parquet(f"{PATH_CLEANED}/text/")
)

# -------------------------------------------------------------------------
# Helper: retry with exponential backoff for throttling (HTTP 429)
# -------------------------------------------------------------------------
def retry_on_throttle(func):
    """Decorator to retry Azure SDK calls on 429 errors with exponential backoff."""
    def wrapper(*args, **kwargs):
        max_retries = 8
        backoff = 5  # start with 5 seconds
        for attempt in range(max_retries):
            try:
                return func(*args, **kwargs)
            except HttpResponseError as e:
                if e.status_code == 429:
                    if attempt == max_retries - 1:
                        raise
                    time.sleep(backoff)
                    backoff = min(backoff * 2, 120)  # cap backoff at 120 sec
                else:
                    raise
    return wrapper

@retry_on_throttle
def _analyze_image_ocr(image_bytes: bytes) -> object:
    """Call Azure Document Intelligence OCR and return the result."""
    client = DocumentAnalysisClient(
        endpoint=DOCUMENTINTEL_ENDPOINT,
        credential=AzureKeyCredential(DOCUMENTINTEL_KEY)
    )
    try:
        poller = client.begin_analyze_document("prebuilt-read", document=io.BytesIO(image_bytes))
        result = poller.result()
        return result
    finally:
        # Explicitly dispose of the client (no close method, but delete reference)
        del client

def redact_image_with_azure_ocr(image_bytes, doc_path):
    log_entries = []
    try:
        # Azure OCR with retry/backoff
        result = _analyze_image_ocr(image_bytes)

        # Load image via Pillow
        image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        draw = ImageDraw.Draw(image)

        for page in result.pages:
            for line in page.lines:
                line_text = line.content
                # Detect language for the line
                try:
                    lang = detect(line_text)
                except Exception:
                    lang = "en"
                # Run Presidio PII detection
                pii_results = analyzer.analyze(text=line_text, language=lang, entities=None)
                # Log line_text and pii_results
                log_entries.append(json.dumps({
                    "doc_path": doc_path,
                    "line_text": line_text,
                    "pii_results": [r.to_dict() for r in pii_results] if pii_results else []
                }, ensure_ascii=False))
                if pii_results:
                    bbox = line.bounding_polygon
                    if bbox:
                        x_coords = [p.x for p in bbox]
                        y_coords = [p.y for p in bbox]
                        min_x, max_x = min(x_coords), max(x_coords)
                        min_y, max_y = min(y_coords), max(y_coords)
                        draw.rectangle([min_x, min_y, max_x, max_y], fill="black")

        # Save redacted image to bytes
        output = io.BytesIO()
        image.save(output, format="PNG")
        redacted_bytes = output.getvalue()

        # Clean up Pillow resources
        draw = None
        image.close()
        output.close()

        # Write log to workspace file
        log_path = "/Workspace/piiresults.txt"
        with open("/tmp/piiresults.txt", "a", encoding="utf-8") as f:
            for entry in log_entries:
                f.write(entry + "\n")
        dbutils.fs.cp("file:/tmp/piiresults.txt", f"file:{log_path}", True)

        return redacted_bytes
    except Exception:
        # In case of any failure, return original bytes
        return image_bytes
    try:
        # Azure OCR with retry/backoff
        result = _analyze_image_ocr(image_bytes)

        # Load image via Pillow
        image = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        draw = ImageDraw.Draw(image)

        for page in result.pages:
            for line in page.lines:
                line_text = line.content
                # Detect language for the line
                try:
                    lang = detect(line_text)
                except Exception:
                    lang = "en"
                # Run Presidio PII detection
                pii_results = analyzer.analyze(text=line_text, language=lang, entities=None)
                if pii_results:
                    bbox = line.bounding_polygon
                    if bbox:
                        x_coords = [p.x for p in bbox]
                        y_coords = [p.y for p in bbox]
                        min_x, max_x = min(x_coords), max(x_coords)
                        min_y, max_y = min(y_coords), max(y_coords)
                        draw.rectangle([min_x, min_y, max_x, max_y], fill="black")

        # Save redacted image to bytes
        output = io.BytesIO()
        image.save(output, format="PNG")
        redacted_bytes = output.getvalue()

        # Clean up Pillow resources
        draw = None
        image.close()
        output.close()

        return redacted_bytes
    except Exception:
        # In case of any failure, return original bytes
        return image_bytes

def list_all_image_files(base_path):
    image_files = []
    def _recurse(path):
        try:
            for entry in dbutils.fs.ls(path):
                if entry.isDir():
                    _recurse(entry.path)
                elif entry.path.lower().endswith((".png", ".jpg", ".jpeg")):
                    image_files.append(entry.path)
        except Exception:
            pass
    _recurse(base_path)
    return image_files

def process_and_redact_images():
    image_files = list_all_image_files(PATH_IMAGES_OUT)
    for img_path in image_files:
        # Copy file to local temp location
        dbutils.fs.cp(img_path, "file:/tmp/tmp_img.png", True)
        with open("/tmp/tmp_img.png", "rb") as f:
            image_bytes = f.read()

        redacted_bytes = redact_image_with_azure_ocr(image_bytes, img_path)

        # Build cleaned path preserving folder hierarchy
        parts = img_path.replace(PATH_IMAGES_OUT, "").lstrip("/").split("/")
        if len(parts) >= 2:
            doc_folder = parts[0]
            sub_path = "/".join(parts[1:])
            cleaned_path = f"{PATH_CLEANED}/images/{doc_folder}/{sub_path}"
        else:
            cleaned_path = f"{PATH_CLEANED}/images/{os.path.basename(img_path)}"

        cleaned_folder = os.path.dirname(cleaned_path)
        dbutils.fs.mkdirs(cleaned_folder)

        # Write redacted image back to ADLS
        with open("/tmp/tmp_img_redacted.png", "wb") as f:
            f.write(redacted_bytes)
        dbutils.fs.cp("file:/tmp/tmp_img_redacted.png", cleaned_path, True)

process_and_redact_images()

# COMMAND ----------

# MAGIC %pip install transformers
# MAGIC # Install PyTorch and Transformers
# MAGIC %pip install torch
# MAGIC %pip install 'accelerate>=0.26.0'
# MAGIC %pip install -U bitsandbytes
# MAGIC dbutils.library.restartPython()

# COMMAND ----------

# ==========================================================
# 03_semantic_problem_solution_classifier_native.py
# ==========================================================
# Purpose:
#   - Read multilingual text documents from ADLS (Parquet)
#   - Compute semantic "problem + solution" scores (no pandas)
#   - Preserve all original columns
#   - Write results back to ADLS as Parquet, partitioned by source folder
# ==========================================================

# COMMAND ----------
from transformers import AutoTokenizer, AutoModel
import torch
from pyspark.sql import functions as F
from pyspark.sql.types import StructType, StructField, StringType, DoubleType
from pyspark import StorageLevel
from pyspark.sql import Row

# ---------- Your provided config ----------
account = "docclassifierstoragecct"
container_raw         = "raw"
container_model       = "docintclassificationmodel"
container_stage       = "stage"
container_redacted    = "redacted"
container_processed   = "docclassifiercontainer"
container_classified  = "classified"

abfss = lambda container, path="": f"abfss://{container}@{account}.dfs.core.windows.net/{path}"

PATH_DOCS_IN     = abfss(container_raw,        "incoming/docs/")      # raw docs (binary) - FYI
PATH_IMAGES_OUT  = abfss(container_stage,      "images/")
PATH_JSON_OUT    = abfss(container_stage,      "extracted_json/")
PATH_CLEANED     = abfss(container_redacted,   "cleaned/text/")            # PII-clean text (if used)


PATH_CLASSIFIED  = abfss(container_classified, "problem_solution/")
LOCAL_MODEL_PATH = "/dbfs/tmp/hf_models/intfloat_multilingual_e5_base"  # quantized model path
MODEL_NAME = LOCAL_MODEL_PATH
SCORE_THRESHOLD = 0.70                          # tune as needed
BATCH_SIZE = 8                                   # per-partition batch for speed
MAX_LENGTH = 512                                 # tokenizer max length per chunk
CHUNK_CHARS = 1500                               # text chunk size in characters

# --- Scan all parquet files in PATH_CLEANED recursively ---
def list_all_parquet_files(base_path):
    files = []
    def _recurse(path):
        try:
            for entry in dbutils.fs.ls(path):
                if entry.isDir():
                    _recurse(entry.path)
                elif entry.path.lower().endswith(".parquet"):
                    files.append(entry.path)
        except Exception:
            pass
    _recurse(base_path)
    return files

all_cleaned_parquet_files = list_all_parquet_files(PATH_CLEANED)

from transformers import AutoTokenizer, AutoModel
import torch

# Load E5 model and tokenizer once on driver
_tokenizer = AutoTokenizer.from_pretrained(LOCAL_MODEL_PATH)
_model = AutoModel.from_pretrained(LOCAL_MODEL_PATH)
_model.eval()

def _embed_text(text: str):
    encoded = _tokenizer([f"passage: {text}"], return_tensors="pt", truncation=True, padding=True, max_length=MAX_LENGTH)
    with torch.no_grad():
        out = _model(**encoded)
    return out.last_hidden_state.mean(dim=1)[0]

emb_problem = _embed_text("This document describes a tax problem or issue.")
emb_solution = _embed_text("This document provides a solution or resolution to a tax problem.")
emb_tax = _embed_text("This document discusses a tax matter.")

def file_contains_tax_matter(parquet_path):
    df = spark.read.parquet(parquet_path)
    if "redacted_text" in df.columns:
        texts = [r["redacted_text"] for r in df.select("redacted_text").limit(1000).collect() if r["redacted_text"]]
    elif "chunk_text" in df.columns:
        texts = [r["chunk_text"] for r in df.select("chunk_text").limit(1000).collect() if r["chunk_text"]]
    else:
        return False
    for text in texts:
        if not text or not isinstance(text, str):
            continue
        # Chunk text for coverage
        chunks = [text[i:i+CHUNK_CHARS] for i in range(0, len(text), CHUNK_CHARS)]
        for chunk in chunks:
            encoded = _tokenizer([f"passage: {chunk}"], return_tensors="pt", truncation=True, padding=True, max_length=MAX_LENGTH)
            with torch.no_grad():
                out = _model(**encoded)
            emb = out.last_hidden_state.mean(dim=1)[0]
            cos = torch.nn.functional.cosine_similarity
            p = cos(emb, emb_problem, dim=0).item()
            s = cos(emb, emb_solution, dim=0).item()
            t = cos(emb, emb_tax, dim=0).item()
            score = max(p, s, t)
            if score >= SCORE_THRESHOLD:
                return True
    return False

import shutil

for parquet_path in all_cleaned_parquet_files:
    if file_contains_tax_matter(parquet_path):
        rel_path = parquet_path.replace(PATH_CLEANED, "").lstrip("/")
        dest_path = f"{PATH_CLASSIFIED}/{rel_path}"
        dest_dir = "/".join(dest_path.split("/")[:-1])
        dbutils.fs.mkdirs(dest_dir)
        dbutils.fs.cp(parquet_path, dest_path, True)
PATH_EMBED       = abfss(container_stage,      "embeddings/")
PATH_TEXT_OUT    = abfss(container_stage,      "text/")               # <-- input parquet (expects columns incl. 'path','text')

# Output for classified docs
PATH_CLASSIFIED  = abfss(container_classified, "problem_solution/")

# ---------- Model / classification settings ----------
LOCAL_MODEL_PATH = "/dbfs/tmp/hf_models/intfloat_multilingual_e5_base"  # quantized model path
MODEL_NAME = LOCAL_MODEL_PATH
SCORE_THRESHOLD = 0.70                          # tune as needed
BATCH_SIZE = 8                                   # per-partition batch for speed
MAX_LENGTH = 512                                 # tokenizer max length per chunk
CHUNK_CHARS = 1500                               # text chunk size in characters

# COMMAND ----------
# ---------- Load model ONCE on driver to build reference embeddings ----------
_tokenizer = AutoTokenizer.from_pretrained(MODEL_NAME)
_model = AutoModel.from_pretrained(MODEL_NAME)
_model.eval()

def _embed_text_driver(text: str):
    """Driver-side helper to embed a single small text with mean pooling."""
    encoded = _tokenizer([f"passage: {text}"], return_tensors="pt", truncation=True, padding=True, max_length=MAX_LENGTH)
    with torch.no_grad():
        out = _model(**encoded)
    # mean pooling over sequence dimension
    return out.last_hidden_state.mean(dim=1)[0]

emb_problem_drv  = _embed_text_driver("This document describes a problem or issue.")
emb_solution_drv = _embed_text_driver("This document provides a solution or resolution.")

# Broadcast references & settings
bc_problem   = spark.sparkContext.broadcast(emb_problem_drv.numpy())
bc_solution  = spark.sparkContext.broadcast(emb_solution_drv.numpy())
bc_modelname = spark.sparkContext.broadcast(MODEL_NAME)
bc_params    = spark.sparkContext.broadcast({
    "batch_size": BATCH_SIZE,
    "max_length": MAX_LENGTH,
    "chunk_chars": CHUNK_CHARS,
    "threshold": SCORE_THRESHOLD
})

# COMMAND ----------
# ---------- Read input parquet (expects at least 'path' and 'text') ----------
# Keep ALL columns; we’ll add to them later.
df_in = spark.read.parquet(PATH_TEXT_OUT)

required_cols = {"path", "chunk_text"}

# Validate essential columns exist
#required_cols = {"path", "text"}
#missing = required_cols - set(df_in.columns)
#if missing:
#    raise ValueError(f"Missing required columns in input parquet at {PATH_TEXT_OUT}: {missing}")

# Persist for reuse
df_in = df_in.filter(F.col("chunk_text").isNotNull())
df_in.persist(StorageLevel.MEMORY_AND_DISK)

# COMMAND ----------
# ---------- Partition-level inference without pandas ----------
def process_partition(rows_iter):
    """
    Spark will call this ONCE per partition.
    We load the model once per partition, then batch process rows.
    """
    import torch
    from transformers import AutoTokenizer, AutoModel

    model_name = bc_modelname.value
    params = bc_params.value
    B = int(params["batch_size"])
    MAXLEN = int(params["max_length"])
    CHUNK = int(params["chunk_chars"])
    TH = float(params["threshold"])

    # Load model/tokenizer inside the executor
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    model = AutoModel.from_pretrained(model_name)
    model.eval()

    emb_problem = torch.tensor(bc_problem.value)
    emb_solution = torch.tensor(bc_solution.value)

    # Helper: chunk chunk_text by characters to improve coverage
    def chunk_text(t: str, n: int):
        t = t or ""
        if len(t) <= n:
            return [t]
        return [t[i:i+n] for i in range(0, len(t), n)]

    # Helper: embed a list of strings with batching
    def embed_texts(text_list):
        embs = []
        # prefix "passage: " for E5
        text_list = [f"passage: {t}" for t in text_list]
        for i in range(0, len(text_list), B):
            batch = text_list[i:i+B]
            enc = tokenizer(batch, return_tensors="pt", padding=True, truncation=True, max_length=MAXLEN)
            with torch.no_grad():
                out = model(**enc)
                # mean pooling
                emb_batch = out.last_hidden_state.mean(dim=1)  # [batch, hidden]
                embs.extend([e for e in emb_batch])
        return embs

    cos = torch.nn.functional.cosine_similarity

    # We’ll carry forward ALL original columns and append new ones.
    results = []
    for row in rows_iter:
        # row is a pyspark Row; convert to dict to extend
        rd = row.asDict(recursive=True)

        text = rd.get("chunk_text")
        if not text or not isinstance(text, str):
            rd["semantic_score"] = 0.0
            rd["classification"] = "Other"
            results.append(Row(**rd))
            continue

        try:
            # chunk & embed
            chunks = chunk_text(text, CHUNK)
            emb_chunks = embed_texts(chunks)

            # compute score over chunks (use MAX for “exists somewhere in doc”)
            best_score = 0.0
            for emb in emb_chunks:
                p = cos(emb, emb_problem, dim=0).item()
                s = cos(emb, emb_solution, dim=0).item()
                score = (p + s) / 2.0
                if score > best_score:
                    best_score = score

            rd["semantic_score"] = float(round(best_score, 4))
            rd["classification"] = "Problem+Solution" if best_score >= TH else "Other"

        except Exception:
            rd["semantic_score"] = 0.0
            rd["classification"] = "Other"

        results.append(Row(**rd))

    return iter(results)

# Run inference
df_scored = spark.createDataFrame(df_in.rdd.mapPartitions(process_partition), schema=df_in.schema.add("semantic_score", "double").add("classification", "string"))

# COMMAND ----------
# ---------- Derive relative folder to preserve source hierarchy ----------
# We try to extract the subpath relative to the RAW incoming/docs/ root (if present in path),
# otherwise fall back to the directory of the provided 'path'.
def _rel_folder_from_path(p: str) -> str:
    """
    Returns a relative folder for partitioning.
    Strategy:
      1) If '/incoming/docs/' is in the path, return the subpath directories after it (without filename).
      2) Else return the immediate directory of the file path.
    """
    if not p:
        return ""
    # Normalize
    s = p.replace("\\", "/")
    # If we see raw incoming/docs path in the full URL, cut after it
    key = "/incoming/docs/"
    if key in s:
        tail = s.split(key, 1)[-1]  # everything after incoming/docs/
        # drop filename if present
        return "/".join(tail.split("/")[:-1])
    # Fallback: directory of file
    return "/".join(s.split("/")[:-1])

extract_rel_udf = F.udf(_rel_folder_from_path, StringType())

df_final = (
    df_scored
    .withColumn("relative_folder", extract_rel_udf(F.col("path")))
)

# COMMAND ----------
# ---------- Write results back to ADLS as Parquet, partitioned by relative_folder ----------
# Adjust coalesce depending on volume; 200 is a default for medium to large jobs.
(
    df_final
    .coalesce(200)
    .write
    .mode("overwrite")
    .partitionBy("relative_folder")
    .parquet(PATH_CLASSIFIED)
)

# COMMAND ----------
# ---------- Quick verification ----------
print("Wrote classified outputs to:", PATH_CLASSIFIED)
display(dbutils.fs.ls(PATH_CLASSIFIED))

# ===========================
# WHY IS THIS SUPERSLOW?
# ===========================
# 1. HuggingFace models (AutoModel) are loaded ONCE PER PARTITION, PER EXECUTOR, from disk.
#    - This is very slow if the model is large or if the model path is on DBFS or remote storage.
#    - Model loading is not parallelized and is repeated for every partition.
# 2. No GPU acceleration: Unless your cluster has GPUs and your Spark workers are configured to use them,
#    inference will run on CPU, which is much slower for transformer models.
# 3. Each row (document chunk) is processed sequentially within the partition, and embedding computation is expensive.
# 4. Large batch sizes or large text chunks can cause memory pressure and slow down processing.
# 5. If the number of partitions is high, model load overhead is multiplied.
# 6. If the model is not cached locally on each worker node, it is reloaded from remote storage every time.
# 7. Spark overhead: Spark is not optimized for deep learning inference workloads.
# 8. Writing with .coalesce(200) can cause shuffling and slow down the write phase if the data is not large enough.

# ===========================
# HOW TO MAKE IT FASTER?
# ===========================
# - Use GPU clusters and ensure torch uses CUDA.
# - Reduce number of partitions to match number of executors.
# - Pre-cache the model on local disk of each worker.
# - Use a lighter model or quantized model.
# - Increase batch size if memory allows.
# - Consider using Databricks Model Serving or a distributed inference service.
# - For large scale, run inference outside Spark (e.g., batch process with Ray, Dask, or TorchServe).

# COMMAND ----------

from pyspark.sql import functions as F, types as T
from transformers import AutoTokenizer, AutoModel
import torch
import os
import re

# Output for classified docs
PATH_CLASSIFIED = abfss(container_classified, "problem_solution/")
PATH_EMBED = abfss(container_stage, "embeddings/")

def list_all_parquet_files(base_path):
    files = []
    def _recurse(path):
        try:
            for entry in dbutils.fs.ls(path):
                if entry.isDir():
                    _recurse(entry.path)
                elif entry.path.lower().endswith(".parquet"):
                    files.append(entry.path)
        except Exception:
            pass
    _recurse(base_path)
    return files

all_parquet_files = list_all_parquet_files(PATH_CLASSIFIED)
print(f"Found {len(all_parquet_files)} parquet files in {PATH_CLASSIFIED}")

LOCAL_MODEL_PATH = "/dbfs/tmp/hf_models/intfloat_multilingual_e5_base"
print(f"Loading model from {LOCAL_MODEL_PATH}")
_tokenizer = AutoTokenizer.from_pretrained(LOCAL_MODEL_PATH)
_model = AutoModel.from_pretrained(LOCAL_MODEL_PATH)
_model.eval()
print("Model loaded successfully.")

def e5_embed_batch(texts):
    if not texts:
        return []
    print(f"Embedding batch of {len(texts)} texts")
    encoded = _tokenizer([f"passage: {t}" if t else "" for t in texts], return_tensors="pt", truncation=True, padding=True, max_length=512)
    with torch.no_grad():
        out = _model(**encoded)
    embs = out.last_hidden_state.mean(dim=1)
    return [e.cpu().numpy().tolist() for e in embs]

def make_doc_id(path, idx):
    base = os.path.basename(path)
    safe_base = re.sub(r'[^A-Za-z0-9_\-=]', '_', base)
    return f"{safe_base}_{idx}"

make_doc_id_udf = F.udf(make_doc_id, T.StringType())

for parquet_path in all_parquet_files:
    print(f"Processing file: {parquet_path}")
    df = spark.read.parquet(parquet_path)
    # Use the correct column for text input
    text_col = None
    for candidate in ["redacted_text", "chunk_text", "content_text"]:
        if candidate in df.columns:
            text_col = candidate
            break
    if not text_col:
        print(f"Skipping {parquet_path}: no valid text column found.")
        continue
    df = df.withColumn("row_idx", F.monotonically_increasing_id())
    texts = [r[text_col] for r in df.select(text_col).collect()]
    print(f"Number of text rows to embed: {len(texts)}")
    if not texts:
        print(f"No texts found in {parquet_path}, skipping.")
        continue
    embeddings = e5_embed_batch(texts)
    print(f"Generated {len(embeddings)} embeddings.")
    emb_rows = [(float(i), emb) for i, emb in enumerate(embeddings)]
    emb_df = spark.createDataFrame(emb_rows, ["row_idx", "embedding"])
    from pyspark.sql.window import Window
    df = df.withColumn("row_idx", F.row_number().over(Window.orderBy(F.monotonically_increasing_id())) - 1)
    df = df.join(emb_df, on="row_idx").drop("row_idx")
    df = df.withColumn("doc_id", make_doc_id_udf(F.col("path"), F.lit(0)))
    select_cols = ["doc_id", "path", "embedding", text_col]
    if "classification" in df.columns:
        select_cols.insert(2, "classification")
    df_out = df.select(*select_cols)
    out_path = f"{PATH_EMBED}/{os.path.basename(parquet_path)}"
    print(f"Writing embeddings to {out_path}")
    df_out.write.mode("overwrite").parquet(out_path)
    print(f"Finished writing {out_path}")

# COMMAND ----------

# Cleaned content to embedding text-embedding-3-large
import os
import io
import json
import re
import requests
import time
from pyspark.sql import functions as F
from pyspark.sql.types import StringType, ArrayType, FloatType
from azure.core.credentials import AzureKeyCredential
from azure.search.documents import SearchClient

# --- List all non-empty parquet files under PATH_TEXT_OUT recursively ---
def list_all_nonempty_parquet_files(base_path):
    files = []
    def _recurse(path):
        try:
            for entry in dbutils.fs.ls(path):
                if entry.isDir():
                    _recurse(entry.path)
                elif entry.path.lower().endswith(".parquet") and entry.size > 0:
                    files.append(entry.path)
        except Exception:
            pass
    _recurse(base_path)
    return files

all_parquet_files = list_all_nonempty_parquet_files(PATH_TEXT_OUT)
if not all_parquet_files:
    raise RuntimeError(f"No non-empty Parquet files found under {PATH_TEXT_OUT}. Aborting embedding pipeline.")

# --- Read all text parquet files (recursively) ---
df_text = spark.read.parquet(*all_parquet_files)

display(df_text.select("path", "chunk_text").limit(10))

def chunk_text(text, chunk_size=1000):
    if not text:
        return []
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

chunk_text_udf = F.udf(chunk_text, ArrayType(StringType()))
df_chunks = df_text.withColumn("chunks", chunk_text_udf("chunk_text"))
df_chunks = df_chunks.select("path", F.posexplode("chunks").alias("chunk_index", "chunk_text"))

display(df_chunks.limit(10))

def retry_on_throttle(func):
    def wrapper(*args, **kwargs):
        max_retries = 8
        backoff = 5
        for attempt in range(max_retries):
            try:
                return func(*args, **kwargs)
            except Exception as e:
                if hasattr(e, "response") and getattr(e.response, "status_code", None) == 429:
                    if attempt == max_retries - 1:
                        raise
                    time.sleep(backoff)
                    backoff = min(backoff * 2, 120)
                elif isinstance(e, requests.exceptions.HTTPError) and e.response.status_code == 429:
                    if attempt == max_retries - 1:
                        raise
                    time.sleep(backoff)
                    backoff = min(backoff * 2, 120)
                elif isinstance(e, requests.exceptions.RequestException) and hasattr(e.response, "status_code") and e.response.status_code == 429:
                    if attempt == max_retries - 1:
                        raise
                    time.sleep(backoff)
                    backoff = min(backoff * 2, 120)
                else:
                    raise
    return wrapper

@retry_on_throttle
def get_embedding(text):
    headers = {
        "Content-Type": "application/json",
        "api-key": AZURE_OPENAI_API_KEY
    }
    data = {
        "input": text,
        "model": "text-embedding-3-large"
    }
    with requests.Session() as session:
        response = session.post(AZURE_OPENAI_EMBEDDING_ENDPOINT, headers=headers, json=data)
        if response.status_code == 200:
            return response.json()["data"][0]["embedding"]
        elif response.status_code == 429:
            raise requests.exceptions.HTTPError(response=response)
        else:
            return []

get_embedding_udf = F.udf(get_embedding, ArrayType(FloatType()))
df_embedded = df_chunks.withColumn("embedding", get_embedding_udf("chunk_text"))

@retry_on_throttle
def push_to_azure_search(docs):
    headers = {
        "Content-Type": "application/json",
        "api-key": AZURE_SEARCH_KEY
    }
    url = f"{AZURE_SEARCH_ENDPOINT}/indexes/{AZURE_SEARCH_INDEX}/docs/index?api-version=2023-11-01"
    with requests.Session() as session:
        response = session.post(url, headers=headers, data=json.dumps({"value": docs}))
        if response.status_code not in (200, 201):
            if response.status_code == 429:
                raise requests.exceptions.HTTPError(response=response)
            raise Exception(f"Failed to push to Azure Search: {response.text}")

def prepare_doc(row):
    base = os.path.basename(row['path'])
    safe_base = re.sub(r'[^A-Za-z0-9_\-=]', '_', base)
    return {
        "id": f"{safe_base}_{row['chunk_index']}",
        "path": row["path"],
        "chunk_index": row["chunk_index"],
        "content": row["chunk_text"],
        "embedding": row["embedding"]
    }

docs = (
    df_embedded
    .select("path", "chunk_index", "chunk_text", "embedding")
    .toLocalIterator()
)

batch = []
batch_size = 100
for row in docs:
    doc = prepare_doc(row.asDict())
    batch.append(doc)
    if len(batch) >= batch_size:
        push_to_azure_search(batch)
        batch = []
if batch:
    push_to_azure_search(batch)

emb = get_embedding("test embedding dimension")
print(len(emb))


# COMMAND ----------

AZURE_SEARCH_ENDPOINT = "https://docsclassifieraisearch.search.windows.net"
AZURE_SEARCH_KEY      = ""
AZURE_SEARCH_INDEX    = "mydocs-knowledgeharvester-index"

AZURE_OPENAI_ENDPOINT = "https://tpapp.openai.azure.com/"
AZURE_OPENAI_KEY      = ""
AZURE_OPENAI_API_VER  = "2024-08-01-preview"
AZURE_OPENAI_MODEL    = "gpt-4o-mini"
AZURE_EMBED_MODEL     = "text-embedding-3-large"

from azure.search.documents import SearchClient
from azure.core.credentials import AzureKeyCredential
from openai import AzureOpenAI

search_client = SearchClient(
    endpoint=AZURE_SEARCH_ENDPOINT,
    index_name=AZURE_SEARCH_INDEX,
    credential=AzureKeyCredential(AZURE_SEARCH_KEY)
)

aoai = AzureOpenAI(
    azure_endpoint=AZURE_OPENAI_ENDPOINT,
    api_key=AZURE_OPENAI_KEY,
    api_version=AZURE_OPENAI_API_VER
)

response = aoai.embeddings.create(
    model="text-embedding-3-large",
    input="test"
)
print(len(response.data[0].embedding))



# COMMAND ----------

# Example usage: Query your indexed documents
response = simple_agent("What is the the information upi have about Lyft?")
print(response)

# COMMAND ----------

print(len(df_embedded.select("embedding").first()["embedding"]))


# COMMAND ----------

# =====================================================
# Agent Configuration
# =====================================================
AZURE_SEARCH_ENDPOINT = "https://docsclassifieraisearch.search.windows.net"
AZURE_SEARCH_KEY      = ""
AZURE_SEARCH_INDEX    = "mydocs-knowledgeharvester-index"

AZURE_OPENAI_ENDPOINT = "https://tpapp.openai.azure.com/"
AZURE_OPENAI_KEY      = ""
AZURE_OPENAI_API_VER  = "2024-08-01-preview"
AZURE_OPENAI_MODEL    = "gpt-4o-mini"
AZURE_EMBED_MODEL     = "text-embedding-3-large"

# =====================================================
# IMPORTS
# =====================================================
from azure.search.documents import SearchClient
from azure.core.credentials import AzureKeyCredential
from openai import AzureOpenAI
import textwrap

# =====================================================
# CLIENTS
# =====================================================
search_client = SearchClient(
    endpoint=AZURE_SEARCH_ENDPOINT,
    index_name=AZURE_SEARCH_INDEX,
    credential=AzureKeyCredential(AZURE_SEARCH_KEY)
)

aoai = AzureOpenAI(
    azure_endpoint=AZURE_OPENAI_ENDPOINT,
    api_key=AZURE_OPENAI_KEY,
    api_version=AZURE_OPENAI_API_VER
)

# =====================================================
# AGENT FUNCTION
# =====================================================
def tax_expert_agent(query: str, top_k: int = 5) -> str:
    """
    1️ Creates embedding for the query
    2️ Performs hybrid retrieval from Azure AI Search
    3️ Uses GPT-4o-mini to rerank and synthesize an answer
    """

    # -------------------------------------------------
    # STEP 1: Get embedding
    # -------------------------------------------------
    embedding = aoai.embeddings.create(
        model=AZURE_EMBED_MODEL,
        input=query
    ).data[0].embedding

    # -------------------------------------------------
    # STEP 2: Hybrid retrieval
    # -------------------------------------------------
    results = search_client.search(
        search_text=query,
        vector_queries=[{
            "kind": "vector",
            "vector": embedding,
            "fields": "embedding",    # ✅ correct field name
            "k": top_k
        }],
        top=top_k,
        query_type="simple"
    )

    docs = list(results)
    if not docs:
        return "Knowledge not found in repository. No web search attempted."

    # -------------------------------------------------
    # STEP 3: Prepare retrieved context for reranking
    # -------------------------------------------------
    passages = []
    for i, d in enumerate(docs, start=1):
        passages.append(f"Document {i}:\n{d.get('content', '')[:2000]}")

    joined_context = "\n\n".join(passages)

    # -------------------------------------------------
    # STEP 4: Build structured prompt for reranking + synthesis
    # -------------------------------------------------
    prompt = f"""
You are a highly professional and courteous **Tax Expert Assistant**.
Use only the knowledge base excerpts provided. Be precise, factual and precise.

If the answer cannot be derived from the knowledge base, say exactly:
"Knowledge not found in repository. No web search attempted."

User Problem Statement:
{query}

Knowledge Base Excerpts (retrieved documents):
{joined_context}

Now:
1. Rerank the above documents based on their factual relevance to the question.
2. Synthesize a clear, concise and polite response.

Format the final output exactly as:

Problem Statement:
Solution:
Additional Information:

Solution can be summarized as bullet points if needed with high level of details as needed.
Never reveal any PII information which identifies any individual.
    """

    # -------------------------------------------------
    # STEP 5: Generate response with GPT-4o-mini
    # -------------------------------------------------
    completion = aoai.chat.completions.create(
        model=AZURE_OPENAI_MODEL,
        messages=[
            {"role": "system", "content": "You are a polite and factual Tax Expert Assistant."},
            {"role": "user", "content": textwrap.dedent(prompt)}
        ],
        temperature=0.2,
        max_tokens=800
    )

    return completion.choices[0].message.content.strip()



# COMMAND ----------

response = tax_expert_agent("Netherlands?")
print(response)